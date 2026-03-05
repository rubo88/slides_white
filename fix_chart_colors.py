"""
Script to convert a dark-background PowerPoint presentation to white background.

Autor: Rubén Veiga
Mail: ruben.veiga@bde.es
Version: v0.1

1. Slide master backgrounds: dark blue (124380) → white (FFFFFF)
2. Upper bar rectangle in masters: dark blue (143966) → grey (D9D9D9)
3. All white text (FFFFFF / schemeClr bg1/lt1) in text boxes → black
4. Chart text → black (titles, axis labels, data labels, legends)
5. Grey pattern bars (D9D9D9/BFBFBF dkHorz) → black (000000/595959)
6. White series lines (FFFFFF) → black (000000)

Handles both traditional charts (c:chartSpace) and extended/waterfall charts
(cx:chartSpace / chartEx).
"""

import os
import glob
import shutil
import zipfile
import tempfile
from pptx import Presentation
from lxml import etree

# Namespaces
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
CX_NS = "http://schemas.microsoft.com/office/drawing/2014/chartex"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Color constants
GREY_FG = "D9D9D9"
GREY_BG = "BFBFBF"
NEW_GREY_FG = "000000"
NEW_GREY_BG = "595959"
WHITE = "FFFFFF"
BLACK = "000000"

DARK_BG = "124380"       # Master background (dark blue)
UPPER_BAR_OLD = "143966"  # Upper bar fill (dark blue)
UPPER_BAR_NEW = "E7E7E7"  # Upper bar fill → grey
WHITE_BG = "FFFFFF"        # New master background


# ===================================================================
#  PART 1 — Slide appearance: backgrounds, upper bar, text colors
# ===================================================================

def fix_all_backgrounds(prs):
    """Force white background on ALL masters, layouts, and slides.

    Sets an explicit solid-white bgPr on every element, removing any
    existing bg (bgRef, bgPr with dark fill, etc.).  This guarantees
    white regardless of theme/inheritance chain.
    """
    changed = 0

    for mi, master in enumerate(prs.slide_masters):
        _set_white_bg(master._element)
        changed += 1
        print(f"  Master {mi+1}: background → white")

    for li, layout in enumerate(prs.slide_layouts):
        _set_white_bg(layout._element)
        changed += 1

    for si, slide in enumerate(prs.slides):
        _set_white_bg(slide._element)
        changed += 1

    print(f"  All {len(list(prs.slide_layouts))} layouts: background → white")
    print(f"  All {len(prs.slides)} slides: background → white")
    return changed


def fix_upper_bar(prs):
    """Change the upper bar rectangle in masters from dark blue to grey."""
    changed = 0
    for mi, master in enumerate(prs.slide_masters):
        for shape in master.shapes:
            if not (hasattr(shape, 'top') and shape.top is not None):
                continue
            top_in = shape.top / 914400
            h_in = shape.height / 914400 if shape.height else 0
            # Upper bar: at the very top, thin (< 1in), full width
            if top_in < 0.05 and h_in < 1.0 and h_in > 0.3:
                sp = shape._element
                spPr = sp.find(f"{{{P_NS}}}spPr")
                if spPr is None:
                    continue
                sf = spPr.find(f"{{{A_NS}}}solidFill")
                if sf is None:
                    continue
                srgb = sf.find(f"{{{A_NS}}}srgbClr")
                if srgb is not None:
                    old_val = srgb.get("val")
                    srgb.set("val", UPPER_BAR_NEW)
                    changed += 1
                    print(f"  Master {mi+1}: upper bar {shape.name} {old_val} → {UPPER_BAR_NEW}")
    return changed


def _clean_bg_shapes(spTree, label):
    """Remove or neutralise full-slide background shapes.

    Removes: images that span the full slide.
    Neutralises: rectangles with dark solid fill or gradient fill
                 that span (nearly) the full slide — replace with
                 white solid fill and no gradient.
    """
    FULL_W = 11000000  # threshold EMU (~full slide width)
    FULL_H = 5000000   # threshold EMU (~full slide height)
    changed = 0

    # Collect elements to process (iterate over a copy because we may remove)
    for sp in list(spTree):
        tag = sp.tag.split("}")[-1] if "}" in sp.tag else sp.tag

        # --- picture shapes → remove ---
        if tag == "pic":
            ext = sp.find(f".//{{{P_NS}}}spPr/{{{A_NS}}}xfrm/{{{A_NS}}}ext")
            if ext is None:
                ext = sp.find(f".//{{{A_NS}}}ext")
            if ext is not None:
                cx = int(ext.get("cx", 0))
                cy = int(ext.get("cy", 0))
                if cx > FULL_W and cy > FULL_H:
                    spTree.remove(sp)
                    changed += 1
                    print(f"  {label}: removed background image")
            continue

        # --- auto-shapes → check if full-size with dark/gradient fill ---
        if tag != "sp":
            continue
        spPr = sp.find(f"{{{P_NS}}}spPr")
        if spPr is None:
            continue

        # Check size
        ext = spPr.find(f"{{{A_NS}}}xfrm/{{{A_NS}}}ext")
        if ext is None:
            continue
        cx = int(ext.get("cx", 0))
        cy = int(ext.get("cy", 0))
        if cx < FULL_W or cy < FULL_H:
            continue  # not full-slide

        # Full-size shape — fix dark fills
        grad = spPr.find(f"{{{A_NS}}}gradFill")
        if grad is not None:
            spPr.remove(grad)
            sf = etree.SubElement(spPr, f"{{{A_NS}}}solidFill")
            clr = etree.SubElement(sf, f"{{{A_NS}}}srgbClr")
            clr.set("val", WHITE_BG)
            changed += 1
            print(f"  {label}: gradient fill → white")
            continue

        solid = spPr.find(f"{{{A_NS}}}solidFill")
        if solid is not None:
            srgb = solid.find(f"{{{A_NS}}}srgbClr")
            scheme = solid.find(f"{{{A_NS}}}schemeClr")
            if srgb is not None and srgb.get("val") in (DARK_BG, "143966"):
                srgb.set("val", WHITE_BG)
                changed += 1
                print(f"  {label}: solid fill → white")
            elif scheme is not None and scheme.get("val") in ("bg2", "dk2"):
                for child in list(solid):
                    solid.remove(child)
                clr = etree.SubElement(solid, f"{{{A_NS}}}srgbClr")
                clr.set("val", WHITE_BG)
                changed += 1
                print(f"  {label}: scheme fill → white")

    return changed


def _set_white_bg(element):
    """Replace any background on element's cSld with solid white."""
    cSld = element.find(f"{{{P_NS}}}cSld")
    if cSld is None:
        return
    bg = cSld.find(f"{{{P_NS}}}bg")
    if bg is not None:
        cSld.remove(bg)
    spTree = cSld.find(f"{{{P_NS}}}spTree")
    new_bg = etree.Element(f"{{{P_NS}}}bg")
    new_bgPr = etree.SubElement(new_bg, f"{{{P_NS}}}bgPr")
    new_sf = etree.SubElement(new_bgPr, f"{{{A_NS}}}solidFill")
    new_clr = etree.SubElement(new_sf, f"{{{A_NS}}}srgbClr")
    new_clr.set("val", WHITE_BG)
    etree.SubElement(new_bgPr, f"{{{A_NS}}}effectLst")
    if spTree is not None:
        idx = list(cSld).index(spTree)
        cSld.insert(idx, new_bg)
    else:
        cSld.insert(0, new_bg)


def fix_blue_fills_in_masters_and_layouts(pptx_path):
    """Zip-level pass: remove ALL blue (124380) solid/gradient shape fills
    from every slide master and layout.  Also removes full-slide picture
    shapes (background images) from masters and layouts.

    Operates directly on the zip so nothing is missed by python-pptx.
    Returns the path of the modified file (overwrites in-place).
    """
    BLUE = "124380"

    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
    os.close(tmp_fd)

    modified = {}

    with zipfile.ZipFile(pptx_path, "r") as zin:
        targets = [f for f in zin.namelist()
                   if (f.startswith("ppt/slideMasters/slideMaster") or
                       f.startswith("ppt/slideLayouts/slideLayout"))
                   and f.endswith(".xml")]

        for fname in targets:
            xml_bytes = zin.read(fname)
            root = etree.fromstring(xml_bytes)
            changes = 0

            # 1. Replace all srgbClr=124380 that are inside a fill context
            #    (solidFill or gradFill gs) with FFFFFF
            for srgb in root.iter(f"{{{A_NS}}}srgbClr"):
                if srgb.get("val", "").upper() != BLUE:
                    continue
                parent = srgb.getparent()
                if parent is None:
                    continue
                ptag = parent.tag.split("}")[-1]
                # Only change fills, not text colours (those are in rPr/defRPr/endParaRPr)
                if ptag in ("solidFill", "gs"):
                    gp = parent.getparent()
                    gptag = gp.tag.split("}")[-1] if gp is not None else ""
                    # Skip if inside a text run property — those are text colours
                    if gptag not in ("rPr", "defRPr", "endParaRPr", "pPr"):
                        srgb.set("val", "FFFFFF")
                        changes += 1

            # 2. Remove full-slide picture shapes (background images)
            for spTree in root.iter(f"{{{P_NS}}}spTree"):
                for pic in list(spTree.findall(f"{{{P_NS}}}pic")):
                    ext = pic.find(f".//{{{A_NS}}}ext")
                    if ext is not None:
                        cx = int(ext.get("cx", 0))
                        cy = int(ext.get("cy", 0))
                        if cx > 10000000 and cy > 5000000:
                            spTree.remove(pic)
                            changes += 1

            if changes:
                modified[fname] = etree.tostring(
                    root, xml_declaration=True, encoding="UTF-8", standalone=True
                )
                print(f"  {fname}: {changes} blue fills/images removed")

    with zipfile.ZipFile(pptx_path, "r") as zin:
        with zipfile.ZipFile(tmp_path, "w") as zout:
            for item in zin.infolist():
                if item.filename in modified:
                    zout.writestr(item, modified[item.filename])
                else:
                    zout.writestr(item, zin.read(item.filename))

    shutil.move(tmp_path, pptx_path)
    return len(modified)


# Scheme color slots that resolve to white on the dark-background theme
WHITE_SCHEME_VALS = {"bg1", "lt1", "tx1", "dk1"}


def _make_black_solidfill():
    """Create a <a:solidFill><a:srgbClr val='000000'/></a:solidFill> element."""
    sf = etree.Element(f"{{{A_NS}}}solidFill")
    clr = etree.SubElement(sf, f"{{{A_NS}}}srgbClr")
    clr.set("val", BLACK)
    return sf


def _force_black_on_element(xml_root):
    """Scan an XML tree and force ALL text to black.

    Replaces or adds explicit black solidFill on every rPr/defRPr/endParaRPr,
    regardless of the current color.

    IMPORTANT: inserts solidFill at position 0 of the rPr element
    so it appears before font elements (latin, ea, cs) — this is
    required by the OOXML schema ordering for PowerPoint to honour it.

    Returns count of elements changed/added.
    """
    changed = 0
    for tag in ["rPr", "defRPr", "endParaRPr"]:
        for rpr in xml_root.iter(f"{{{A_NS}}}{tag}"):
            sf = rpr.find(f"{{{A_NS}}}solidFill")
            if sf is not None:
                rpr.remove(sf)
            rpr.insert(0, _make_black_solidfill())
            changed += 1
    return changed


def _fix_all_borders(xml_root):
    """Change ALL shape borders/outlines to black.

    Finds <a:ln><a:solidFill> and forces the color to black,
    regardless of the current color value.
    """
    changed = 0
    for ln in xml_root.iter(f"{{{A_NS}}}ln"):
        sf = ln.find(f"{{{A_NS}}}solidFill")
        if sf is None:
            continue
        # Replace all fill contents with black srgbClr
        for child in list(sf):
            sf.remove(child)
        clr = etree.SubElement(sf, f"{{{A_NS}}}srgbClr")
        clr.set("val", BLACK)
        changed += 1
    return changed


def fix_all_text_and_borders(prs):
    """Force ALL text across masters, layouts and slides to black.
    Also force ALL borders/outlines on shapes to black.

    Scans the full XML of each element (including txStyles,
    placeholders, tables, groups, SmartArt — everything).
    """
    text_total = 0
    border_total = 0

    # Masters (full _element → includes txStyles, shapes, everything)
    for mi, master in enumerate(prs.slide_masters):
        nt = _force_black_on_element(master._element)
        nb = _fix_all_borders(master._element)
        if nt or nb:
            text_total += nt
            border_total += nb
            parts = []
            if nt:
                parts.append(f"{nt} text")
            if nb:
                parts.append(f"{nb} borders")
            print(f"  Master {mi+1}: {', '.join(parts)} → black")

    # Layouts
    for li, layout in enumerate(prs.slide_layouts):
        nt = _force_black_on_element(layout._element)
        nb = _fix_all_borders(layout._element)
        if nt or nb:
            text_total += nt
            border_total += nb
            parts = []
            if nt:
                parts.append(f"{nt} text")
            if nb:
                parts.append(f"{nb} borders")
            print(f"  Layout {li+1}: {', '.join(parts)} → black")

    # Slides
    for si, slide in enumerate(prs.slides):
        nt = _force_black_on_element(slide._element)
        nb = _fix_all_borders(slide._element)
        if nt or nb:
            text_total += nt
            border_total += nb
            parts = []
            if nt:
                parts.append(f"{nt} text")
            if nb:
                parts.append(f"{nb} borders")
            print(f"  Slide {si+1}: {', '.join(parts)} → black")

    return text_total + border_total


# ===================================================================
#  PART 2 — Chart color fixes
# ===================================================================

def make_black_solid_fill():
    sf = etree.Element(f"{{{A_NS}}}solidFill")
    clr = etree.SubElement(sf, f"{{{A_NS}}}srgbClr")
    clr.set("val", BLACK)
    return sf


def set_text_color_black(rpr):
    for sf in rpr.findall(f"{{{A_NS}}}solidFill"):
        rpr.remove(sf)
    rpr.insert(0, make_black_solid_fill())


def get_srgb_val(color_container):
    if color_container is None:
        return None
    srgb = color_container.find(f"{{{A_NS}}}srgbClr")
    return srgb.get("val") if srgb is not None else None


def set_srgb_val(color_container, new_val):
    srgb = color_container.find(f"{{{A_NS}}}srgbClr")
    if srgb is not None:
        srgb.set("val", new_val)


def is_grey_dkhorz_pattern(spPr):
    patt = spPr.find(f"{{{A_NS}}}pattFill")
    if patt is None or patt.get("prst") != "dkHorz":
        return False
    fg_val = get_srgb_val(patt.find(f"{{{A_NS}}}fgClr"))
    bg_val = get_srgb_val(patt.find(f"{{{A_NS}}}bgClr"))
    return fg_val == GREY_FG and bg_val in (GREY_BG, WHITE)


def fix_grey_pattern(spPr):
    patt = spPr.find(f"{{{A_NS}}}pattFill")
    if patt is None:
        return
    set_srgb_val(patt.find(f"{{{A_NS}}}fgClr"), NEW_GREY_FG)
    set_srgb_val(patt.find(f"{{{A_NS}}}bgClr"), NEW_GREY_BG)


def get_line_srgb_val(spPr):
    ln = spPr.find(f"{{{A_NS}}}ln")
    if ln is None:
        return None
    sf = ln.find(f"{{{A_NS}}}solidFill")
    if sf is None:
        return None
    srgb = sf.find(f"{{{A_NS}}}srgbClr")
    if srgb is not None:
        return srgb.get("val")
    scheme = sf.find(f"{{{A_NS}}}schemeClr")
    if scheme is not None and scheme.get("val") == "bg1":
        return WHITE
    return None


def fix_white_line(spPr):
    ln = spPr.find(f"{{{A_NS}}}ln")
    if ln is None:
        return
    sf = ln.find(f"{{{A_NS}}}solidFill")
    if sf is None:
        return
    srgb = sf.find(f"{{{A_NS}}}srgbClr")
    if srgb is not None and srgb.get("val") == WHITE:
        srgb.set("val", BLACK)
        return
    scheme = sf.find(f"{{{A_NS}}}schemeClr")
    if scheme is not None and scheme.get("val") == "bg1":
        children = list(scheme)
        sf.remove(scheme)
        new_clr = etree.SubElement(sf, f"{{{A_NS}}}srgbClr")
        new_clr.set("val", BLACK)
        for child in children:
            new_clr.append(child)


def process_text_colors(root):
    for tag in ["defRPr", "rPr", "endParaRPr"]:
        for elem in root.iter(f"{{{A_NS}}}{tag}"):
            set_text_color_black(elem)


def process_axis_lines(root):
    for ax_tag in ["catAx", "valAx", "dateAx", "serAx"]:
        for ax in root.iter(f"{{{C_NS}}}{ax_tag}"):
            for spPr in ax.findall(f"{{{C_NS}}}spPr"):
                for ln in spPr.findall(f"{{{A_NS}}}ln"):
                    for sf in ln.findall(f"{{{A_NS}}}solidFill"):
                        for child in list(sf):
                            sf.remove(child)
                        clr = etree.SubElement(sf, f"{{{A_NS}}}srgbClr")
                        clr.set("val", BLACK)


def process_series_fills_traditional(root):
    changed = 0
    for ser in root.iter(f"{{{C_NS}}}ser"):
        spPr = ser.find(f"{{{C_NS}}}spPr")
        if spPr is None:
            continue
        if is_grey_dkhorz_pattern(spPr):
            fix_grey_pattern(spPr)
            changed += 1
        if get_line_srgb_val(spPr) == WHITE:
            fix_white_line(spPr)
            changed += 1
    return changed


def process_series_fills_chartex(root):
    changed = 0
    for dataPt in root.iter(f"{{{CX_NS}}}dataPt"):
        spPr = dataPt.find(f"{{{CX_NS}}}spPr")
        if spPr is None:
            continue
        if is_grey_dkhorz_pattern(spPr):
            fix_grey_pattern(spPr)
            changed += 1
    for series in root.iter(f"{{{CX_NS}}}series"):
        spPr = series.find(f"{{{CX_NS}}}spPr")
        if spPr is None:
            continue
        if is_grey_dkhorz_pattern(spPr):
            fix_grey_pattern(spPr)
            changed += 1
    return changed


def process_traditional_charts(prs):
    count = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_chart:
                cs = shape.chart._chartSpace
                process_text_colors(cs)
                process_axis_lines(cs)
                fill_changes = process_series_fills_traditional(cs)
                count += 1
                title = ""
                if shape.chart.has_title:
                    try:
                        title = shape.chart.chart_title.text_frame.text
                    except Exception:
                        title = "(no text)"
                extra = f" ({fill_changes} fill/line fixes)" if fill_changes else ""
                print(f"  [chart]   Slide {slide_idx + 1}: {title}{extra}")
    return count


def process_chartex_files(pptx_path):
    count = 0
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
    os.close(tmp_fd)
    shutil.copy2(pptx_path, tmp_path)

    chartex_files = []
    with zipfile.ZipFile(pptx_path, "r") as zin:
        chartex_files = [
            f for f in zin.namelist()
            if f.startswith("ppt/charts/chartEx") and f.endswith(".xml")
        ]

    if not chartex_files:
        os.unlink(tmp_path)
        return 0

    modified_contents = {}
    with zipfile.ZipFile(pptx_path, "r") as zin:
        for cex_file in chartex_files:
            xml_bytes = zin.read(cex_file)
            root = etree.fromstring(xml_bytes)
            process_text_colors(root)
            fill_changes = process_series_fills_chartex(root)
            modified_contents[cex_file] = etree.tostring(
                root, xml_declaration=True, encoding="UTF-8", standalone=True
            )
            count += 1
            title = ""
            for t_elem in root.iter(f"{{{A_NS}}}t"):
                title = t_elem.text or ""
                break
            extra = f" ({fill_changes} fill fixes)" if fill_changes else ""
            print(f"  [chartEx] {cex_file}: {title}{extra}")

    with zipfile.ZipFile(pptx_path, "r") as zin:
        with zipfile.ZipFile(tmp_path, "w") as zout:
            for item in zin.infolist():
                if item.filename in modified_contents:
                    zout.writestr(item, modified_contents[item.filename])
                else:
                    zout.writestr(item, zin.read(item.filename))

    shutil.move(tmp_path, pptx_path)
    return count


def fix_theme_color_schemes(pptx_path):
    """Fix the theme dk1/dk2 colors from white → black.

    In the 'Plantilla azul 2026' dark-background themes, dk1 (the default
    text color) is FFFFFF (white).  On a white background this must be
    000000 (black).  We also fix dk2 and accent1 if they are all-white.
    Operates at the zip level (themes are not exposed by python-pptx).
    """
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
    os.close(tmp_fd)

    # Map of theme scheme slot → new srgbClr value
    SLOT_FIXES = {
        "dk1": "000000",   # default text color
        "dk2": "333333",   # secondary text color
    }

    modified_contents = {}
    with zipfile.ZipFile(pptx_path, "r") as zin:
        theme_files = [
            f for f in zin.namelist()
            if f.startswith("ppt/theme/theme") and f.endswith(".xml")
        ]
        for tf in theme_files:
            xml_bytes = zin.read(tf)
            root = etree.fromstring(xml_bytes)
            changed = 0
            for clrScheme in root.iter(f"{{{A_NS}}}clrScheme"):
                for slot, new_val in SLOT_FIXES.items():
                    el = clrScheme.find(f"{{{A_NS}}}{slot}")
                    if el is None:
                        continue
                    srgb = el.find(f"{{{A_NS}}}srgbClr")
                    sys_clr = el.find(f"{{{A_NS}}}sysClr")
                    if srgb is not None and srgb.get("val") == WHITE:
                        srgb.set("val", new_val)
                        changed += 1
                    elif sys_clr is not None and sys_clr.get("lastClr") == WHITE:
                        # Replace sysClr with srgbClr for explicit control
                        el.remove(sys_clr)
                        new_srgb = etree.SubElement(el, f"{{{A_NS}}}srgbClr")
                        new_srgb.set("val", new_val)
                        changed += 1
                # Also fix accent1 if it is pure white (invisible on white bg)
                acc1 = clrScheme.find(f"{{{A_NS}}}accent1")
                if acc1 is not None:
                    srgb = acc1.find(f"{{{A_NS}}}srgbClr")
                    if srgb is not None and srgb.get("val") == WHITE:
                        srgb.set("val", "124380")  # reuse the original blue
                        changed += 1
            if changed:
                modified_contents[tf] = etree.tostring(
                    root, xml_declaration=True, encoding="UTF-8", standalone=True
                )
                print(f"  {tf}: {changed} color slot(s) fixed")

    if not modified_contents:
        return 0

    with zipfile.ZipFile(pptx_path, "r") as zin:
        with zipfile.ZipFile(tmp_path, "w") as zout:
            for item in zin.infolist():
                if item.filename in modified_contents:
                    zout.writestr(item, modified_contents[item.filename])
                else:
                    zout.writestr(item, zin.read(item.filename))

    shutil.move(tmp_path, pptx_path)
    return len(modified_contents)


def fix_text_and_borders_in_zip(pptx_path):
    """Zip-level pass: force ALL text to black and ALL borders to black
    in every slide master and layout XML.

    python-pptx's prs.slide_layouts only returns ~3 of 45 layouts, so
    this function ensures every layout file in the zip is covered.
    """
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
    os.close(tmp_fd)

    modified = {}

    with zipfile.ZipFile(pptx_path, "r") as zin:
        targets = [f for f in zin.namelist()
                   if (f.startswith("ppt/slideMasters/slideMaster") or
                       f.startswith("ppt/slideLayouts/slideLayout"))
                   and f.endswith(".xml")]

        for fname in targets:
            xml_bytes = zin.read(fname)
            root = etree.fromstring(xml_bytes)
            changes = 0

            # Force all text run properties to explicit black
            for tag in ["rPr", "defRPr", "endParaRPr"]:
                for rpr in root.iter(f"{{{A_NS}}}{tag}"):
                    sf = rpr.find(f"{{{A_NS}}}solidFill")
                    if sf is not None:
                        rpr.remove(sf)
                    new_sf = etree.Element(f"{{{A_NS}}}solidFill")
                    clr = etree.SubElement(new_sf, f"{{{A_NS}}}srgbClr")
                    clr.set("val", BLACK)
                    rpr.insert(0, new_sf)
                    changes += 1

            # Force all shape borders to black
            for ln in root.iter(f"{{{A_NS}}}ln"):
                sf = ln.find(f"{{{A_NS}}}solidFill")
                if sf is None:
                    continue
                for child in list(sf):
                    sf.remove(child)
                clr = etree.SubElement(sf, f"{{{A_NS}}}srgbClr")
                clr.set("val", BLACK)
                changes += 1

            if changes:
                modified[fname] = etree.tostring(
                    root, xml_declaration=True, encoding="UTF-8", standalone=True
                )
                print(f"  {fname}: {changes} changes")

    with zipfile.ZipFile(pptx_path, "r") as zin:
        with zipfile.ZipFile(tmp_path, "w") as zout:
            for item in zin.infolist():
                if item.filename in modified:
                    zout.writestr(item, modified[item.filename])
                else:
                    zout.writestr(item, zin.read(item.filename))

    shutil.move(tmp_path, pptx_path)
    return len(modified)


# ===================================================================
#  Main
# ===================================================================

def process_file(input_file, output_file):
    print(f"\nProcessing {input_file} -> {output_file}")
    print("=" * 60)
    print("PART 1: Slide appearance (backgrounds, bar, text)")
    print("=" * 60)

    prs = Presentation(input_file)

    fix_all_backgrounds(prs)
    fix_upper_bar(prs)
    fix_all_text_and_borders(prs)

    print("\n" + "=" * 60)
    print("PART 2: Chart colors")
    print("=" * 60)

    trad_count = process_traditional_charts(prs)
    prs.save(output_file)
    print(f"\n  Traditional charts processed: {trad_count}")

    chartex_count = process_chartex_files(output_file)
    print(f"  Extended (chartEx) charts processed: {chartex_count}")

    print("\n" + "=" * 60)
    print("PART 3: Remove blue fills from masters/layouts")
    print("=" * 60)
    blue_count = fix_blue_fills_in_masters_and_layouts(output_file)
    print(f"  Files with blue fills removed: {blue_count}")

    print("\n" + "=" * 60)
    print("PART 3b: Force text+borders black in ALL masters/layouts (zip)")
    print("=" * 60)
    zip_text_count = fix_text_and_borders_in_zip(output_file)
    print(f"  Files updated: {zip_text_count}")

    print("\n" + "=" * 60)
    print("PART 4: Theme color scheme fixes")
    print("=" * 60)

    theme_count = fix_theme_color_schemes(output_file)
    print(f"  Theme files fixed: {theme_count}")

    print(f"\nDone with {input_file}! Total: {trad_count + chartex_count} charts, {theme_count} themes, {blue_count} blue-fill files.")
    print(f"Saved to: {output_file}")


def main():
    folder_path = os.path.dirname(os.path.abspath(__file__))
    for input_file in glob.glob(os.path.join(folder_path, "*.pptx")):
        basename = os.path.basename(input_file)
        if basename.endswith("_white.pptx") or basename.startswith("~$"):
            continue
        
        base_name_no_ext = os.path.splitext(input_file)[0]
        output_file = f"{base_name_no_ext}_white.pptx"
        process_file(input_file, output_file)

if __name__ == "__main__":
    main()
